"""
main.py — FastAPI Backend (Web Bridge)
----------------------------------------
Connects the React frontend to the LLM planning layer and the dual-MCP
rendering pipeline. Exposes four HTTP endpoints:

  POST /generate      — calls generator.generate_slides() and returns JSON
  POST /export        — builds a .pptx via MCP and streams it as a download
  GET  /history       — returns previously generated presentations from .history/
  GET  /image-proxy   — proxies Pexels image requests to avoid CORS in the browser

Why a separate HTTP server instead of calling MCP directly from the browser?
  - Browsers cannot spawn subprocesses, so the stdio MCP transport must live
    server-side. FastAPI acts as the thin bridge between HTTP and stdio MCP.
  - The Pexels API key must stay server-side to avoid exposing it in the browser.

Run:
  uvicorn backend.main:app --reload --port 8000
"""

import sys
import os
import re
import json
import asyncio
import base64
import io

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response
from pydantic import BaseModel
from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client
from pptx import Presentation
from pptx.util import Inches
import requests as req_lib
from dotenv import load_dotenv

# ── Path & env setup ──────────────────────────────────────────────────────────
# ROOT_MODULE = assignment/ — used to locate .env, docs/, and .history/
ROOT_MODULE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
load_dotenv(os.path.join(ROOT_MODULE, ".env"))

# Import the LLM planning function (shared with the CLI agent)
from .generator import generate_slides

# ── FastAPI app ───────────────────────────────────────────────────────────────
app = FastAPI()

# Allow all origins so the Vite dev server (port 3000) can call this API (port 8000).
# In production this should be restricted to the actual frontend domain.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


# ── Request / Response models ─────────────────────────────────────────────────

class GenerateRequest(BaseModel):
    """Body for POST /generate."""
    topic:        str
    requirements: str = ""   # optional extra instructions from the user
    num_slides:   int = 6    # how many content slides to generate


class ExportRequest(BaseModel):
    """Body for POST /export — the frontend sends back the JSON it received from /generate."""
    presentation: dict        # full presentation dict from /generate
    fetch_images: bool = True # whether to call Pexels for images


class CanvasExportRequest(BaseModel):
    """Body for POST /export-canvas (legacy — kept for compatibility)."""
    slides: list[str]         # list of base64-encoded PNG screenshots
    title:  str = "presentation"


# ── POST /generate ────────────────────────────────────────────────────────────

@app.post("/generate")
async def generate(req: GenerateRequest):
    """
    Plan the full presentation using the LLM and return the structured JSON.

    This is the first step in the web flow:
      Browser → POST /generate → generator.generate_slides() → JSON response

    The returned JSON is displayed in the React SlideCanvas and also sent back
    to POST /export when the user clicks "Export PPTX".

    Side effect: saves the JSON plan to .history/ for the /history endpoint.
    """
    try:
        data = generate_slides(req.topic, req.requirements, req.num_slides)

        # Persist the plan to .history/ so it appears in the history panel
        clean_title = re.sub(r'[^\w\s-]', '', data.get('title', 'presentation')).strip()[:60].replace(' ', '_')
        if not clean_title:
            clean_title = "presentation"

        history_dir = os.path.join(ROOT_MODULE, ".history")
        os.makedirs(history_dir, exist_ok=True)

        # Avoid overwriting existing history entries with the same title
        json_path = os.path.join(history_dir, f"{clean_title}.json")
        seq = 1
        while os.path.exists(json_path):
            json_path = os.path.join(history_dir, f"{clean_title}_{seq}.json")
            seq += 1

        with open(json_path, "w") as f:
            json.dump(data, f)

        return data

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── GET /history ──────────────────────────────────────────────────────────────

@app.get("/history")
async def get_history():
    """
    Return all previously generated presentations, newest first.

    Reads JSON files from .history/ and returns them as a list.
    The frontend uses this to populate the history sidebar panel.
    """
    history_dir = os.path.join(ROOT_MODULE, ".history")
    history = []

    if os.path.exists(history_dir):
        # Collect all .json files and sort by modification time (newest first)
        files = sorted(
            [os.path.join(history_dir, f) for f in os.listdir(history_dir) if f.endswith(".json")],
            key=os.path.getmtime,
            reverse=True,
        )
        for filepath in files:
            with open(filepath, "r") as f:
                try:
                    history.append(json.load(f))
                except Exception:
                    pass  # skip malformed JSON files silently

    return history


# ── POST /export (MCP pipeline) ───────────────────────────────────────────────

async def _build_pptx_via_mcp(presentation: dict, fetch_images: bool) -> bytes:
    """
    Build a .pptx file by orchestrating the two MCP servers over stdio.

    This is the core of the export pipeline:
      1. Spawn mcp_server.py       (PPT rendering — python-pptx)
      2. Spawn web_search_mcp.py   (image search — Pexels API)
      3. Call tools in sequence:
           create_presentation → add_title_slide → add_slide × N → save_presentation
      4. Read the saved .pptx bytes and return them to the caller

    Args:
        presentation: The full presentation dict from /generate.
        fetch_images: Whether to call search_image for each slide.

    Returns:
        Raw .pptx bytes ready to stream back to the browser.
    """
    # Save the .pptx to docs/ as a persistent copy alongside the download
    docs_dir = os.path.join(ROOT_MODULE, "docs")
    os.makedirs(docs_dir, exist_ok=True)
    filename    = f"{presentation.get('title', 'presentation').replace(' ', '_')}.pptx"
    output_path = os.path.join(docs_dir, filename)

    # ── Spawn both MCP servers as subprocesses ────────────────────────────────
    # StdioServerParameters tells the MCP client to launch the server script
    # as a child process and communicate via stdin/stdout pipes.
    server_params = StdioServerParameters(
        command=sys.executable,
        args=[os.path.join(os.path.dirname(__file__), "mcp_server.py")],
    )
    server_params_search = StdioServerParameters(
        command=sys.executable,
        args=[os.path.join(os.path.dirname(__file__), "web_search_mcp.py")],
    )

    async with stdio_client(server_params)        as (read, write), \
               stdio_client(server_params_search) as (read_search, write_search):
        async with ClientSession(read, write)               as session, \
                   ClientSession(read_search, write_search) as session_search:

            # MCP handshake — must complete before any tool calls
            await session.initialize()
            await session_search.initialize()

            # Initialise an empty Presentation object inside mcp_server.py
            await session.call_tool("create_presentation", {"output_path": output_path})

            # ── Title slide ───────────────────────────────────────────────────
            # Fetch a cover image then render the title slide
            cover_query = presentation.get("cover_image_query") or presentation.get("title")
            cover_url   = ""
            if fetch_images:
                r_img     = await session_search.call_tool("search_image", {"query": cover_query})
                cover_url = r_img.content[0].text

            await session.call_tool("add_title_slide", {
                "title":     presentation["title"],
                "subtitle":  presentation.get("subtitle", ""),
                "image_url": cover_url,
            })

            # ── Content slides ────────────────────────────────────────────────
            # Iterate over each planned slide, fetch its image, then render it.
            # All layout fields are passed even if unused — mcp_server.py ignores
            # fields that don't apply to the chosen layout.
            for slide in presentation.get("slides", []):
                slide_url = ""
                if fetch_images and slide.get("image_query"):
                    r_img     = await session_search.call_tool("search_image", {"query": slide["image_query"]})
                    slide_url = r_img.content[0].text

                await session.call_tool("add_slide", {
                    "title":         slide["title"],
                    "layout":        slide.get("layout", "bullets"),
                    "bullets":       slide.get("bullets", []),
                    "left_bullets":  slide.get("left_bullets", []),
                    "right_bullets": slide.get("right_bullets", []),
                    "quote":         slide.get("quote", ""),
                    "author":        slide.get("author", ""),
                    "stats":         slide.get("stats", []),
                    "description":   slide.get("description", ""),
                    "image_url":     slide_url,
                    "notes":         slide.get("notes", ""),
                })

            # Flush the in-memory Presentation to disk — must be called last
            await session.call_tool("save_presentation", {})

    # Read the saved file and return raw bytes to the HTTP response
    with open(output_path, "rb") as f:
        return f.read()


@app.post("/export")
async def export(req: ExportRequest):
    """
    Build and stream a .pptx file matching the UI preview.

    Receives the presentation JSON that the frontend already has from /generate,
    passes it through the MCP pipeline, and returns the .pptx as a file download.

    The MCP server replicates the exact same visual style as the React SlideCard
    components (colors, fonts, layout geometry) so the download matches the preview.
    """
    try:
        pptx_bytes = await _build_pptx_via_mcp(req.presentation, req.fetch_images)
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── POST /export-canvas (legacy) ──────────────────────────────────────────────

@app.post("/export-canvas")
async def export_canvas(req: CanvasExportRequest):
    """
    Legacy endpoint: build a .pptx from browser-captured slide screenshots.

    Each slide is a base64-encoded PNG (captured via html2canvas in the browser).
    This approach is kept for compatibility but is no longer used by the frontend —
    the /export endpoint (MCP pipeline) produces better results because:
      - html2canvas cannot capture cross-origin Pexels images (CORS)
      - Screenshots are taken before React re-renders, producing blank slides

    Args (via request body):
        slides: List of base64 PNG strings, one per slide.
        title:  Presentation title used for the filename.
    """
    try:
        # Build a minimal Presentation with one full-bleed image slide per screenshot
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]  # blank layout — no placeholders

        for b64 in req.slides:
            slide      = prs.slides.add_slide(blank)
            img_bytes  = base64.b64decode(b64)
            # Stretch the screenshot to fill the entire slide
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(0), Inches(0),
                prs.slide_width, prs.slide_height,
            )

        buf = io.BytesIO()
        prs.save(buf)

        # Sanitise the title for use as a filename
        clean_title = re.sub(r'[^\w\s-]', '', req.title).strip()[:60].replace(' ', '_')
        if not clean_title:
            clean_title = "presentation"

        # Save a copy to docs/ and avoid overwriting existing files
        docs_dir  = os.path.join(ROOT_MODULE, "docs")
        os.makedirs(docs_dir, exist_ok=True)
        docs_path = os.path.join(docs_dir, f"{clean_title}.pptx")
        seq = 1
        while os.path.exists(docs_path):
            docs_path = os.path.join(docs_dir, f"{clean_title}_{seq}.pptx")
            seq += 1

        with open(docs_path, "wb") as f:
            f.write(buf.getvalue())

        return Response(
            content=buf.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={os.path.basename(docs_path)}"},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── GET /image-proxy ──────────────────────────────────────────────────────────

@app.get("/image-proxy")
async def image_proxy(query: str):
    """
    Proxy Pexels image requests from the browser to avoid CORS issues.

    The browser's SlideCard components call /image-proxy?query=... to display
    images in the preview. Direct Pexels API calls from the browser would fail
    because the API key must be sent in a header (not a URL param) and the
    Pexels CDN does not set permissive CORS headers for all origins.

    This endpoint:
      1. Searches Pexels for landscape photos matching the query
      2. Downloads the first result's large image
      3. Streams the raw image bytes back to the browser

    Args (query param):
        query: Pexels search string, e.g. "climate change forest".
    """
    api_key = os.getenv("PEXELS_API_KEY", "")
    if not api_key:
        raise HTTPException(status_code=404, detail="No Pexels API key configured")

    try:
        # Search Pexels for landscape photos matching the query
        r = req_lib.get(
            "https://api.pexels.com/v1/search",
            headers={"Authorization": api_key},
            params={"query": query, "per_page": 5, "orientation": "landscape"},
            timeout=10,
        )
        photos = r.json().get("photos", [])
        if not photos:
            raise HTTPException(status_code=404, detail="No images found")

        # Download the first result and stream it back to the browser
        img = req_lib.get(photos[0]["src"]["large"], timeout=10)
        if img.status_code == 200:
            return Response(
                content=img.content,
                media_type=img.headers.get("Content-Type", "image/jpeg"),
            )

    except HTTPException:
        raise  # re-raise our own 404s
    except Exception:
        pass   # network errors fall through to the 404 below

    raise HTTPException(status_code=404, detail="Image not found")
