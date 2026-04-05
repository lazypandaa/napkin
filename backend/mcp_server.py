"""
MCP Server — full napkin_ai-quality PPT tools over stdio
Tools: create_presentation, 
        add_title_slide, 
        add_slide, 
        save_presentation
"""
import io
import os
import sys
import asyncio
from typing import Optional
from dotenv import load_dotenv
import requests
from lxml import etree

load_dotenv(os.path.join(os.path.dirname(__file__), "..", ".env"))

from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp import types
from pptx import Presentation as PRS
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme (matches SlideCard.css) ────────────────────────────────────────────
THEME = {
    "bg":         RGBColor(0xFF, 0xFF, 0xFF),  # white
    "title_color":RGBColor(0x1A, 0x1A, 0x1A),  # #1a1a1a
    "text_color": RGBColor(0x2A, 0x2A, 0x2A),  # #2a2a2a
    "accent":     RGBColor(0x6C, 0x47, 0xFF),  # #6c47ff purple
    "muted":      RGBColor(0x66, 0x66, 0x66),  # #666
    "stat_color": RGBColor(0x6C, 0x47, 0xFF),  # #6c47ff purple
    "stat_bg":    RGBColor(0xF7, 0xF5, 0xFF),  # #f7f5ff
    "col_label":  RGBColor(0x6C, 0x47, 0xFF),  # #6c47ff
}
PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── State ─────────────────────────────────────────────────────────────────────
_prs: Optional[PRS] = None
_output_path: str = "output.pptx"
_used_image_urls: set = set()

# ── Drawing helpers (ported from exporter.py) ─────────────────────────────────

def _set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or THEME["bg"]

def _textbox(slide, text, left, top, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tb

def _divider(slide, top):
    # 48px wide purple divider matching .slide-divider in CSS
    line = slide.shapes.add_shape(1, Inches(0.5), top, Inches(0.75), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = THEME["accent"]
    line.line.fill.background()

def _bullets_block(slide, bullets, left, top, width, height, font_size=16):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"▸  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = THEME["text_color"]

def _description_block(slide, text, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.color.rgb = THEME["muted"]
    run.font.italic = True

def _download_image(url: str) -> Optional[bytes]:
    global _used_image_urls
    if not url or url in _used_image_urls:
        return None
    try:
        img = requests.get(url, timeout=10)
        if img.status_code == 200 and "image" in img.headers.get("Content-Type", ""):
            _used_image_urls.add(url)
            return img.content
    except Exception:
        pass
    return None

def _add_image(slide, url, left, top, width, height) -> bool:
    data = _download_image(url)
    if data:
        try:
            slide.shapes.add_picture(io.BytesIO(data), left, top, width, height)
            return True
        except Exception:
            pass
    return False

def _add_dark_overlay(slide, alpha_pct: int = 55):
    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    overlay.line.fill.background()
    sp = overlay._element
    spPr = sp.find(f"{{{PML}}}spPr")
    for child in list(spPr):
        if etree.QName(child.tag).localname in ("solidFill", "gradFill", "noFill", "pattFill"):
            spPr.remove(child)
    solidFill = etree.SubElement(spPr, f"{{{DML}}}solidFill")
    srgbClr   = etree.SubElement(solidFill, f"{{{DML}}}srgbClr")
    srgbClr.set("val", "000000")
    alpha_el  = etree.SubElement(srgbClr, f"{{{DML}}}alpha")
    alpha_el.set("val", str(alpha_pct * 1000))

# ── Layout renderers (ported from exporter.py) ────────────────────────────────

def _render_bullets(slide, data):
    has_image = bool(data.get("image_url"))
    content_width = Inches(7.5) if has_image else Inches(12.3)
    if data.get("description"):
        _description_block(slide, data["description"], Inches(0.5), Inches(1.55), content_width, Inches(0.8))
    _bullets_block(slide, data.get("bullets", []), Inches(0.5), Inches(2.45), content_width, Inches(4.5))
    if has_image:
        _add_image(slide, data["image_url"], Inches(8.1), Inches(0), Inches(5.23), Inches(7.5))

def _render_two_column(slide, data):
    left_b  = data.get("left_bullets")  or data.get("bullets", [])[:3]
    right_b = data.get("right_bullets") or data.get("bullets", [])[3:]
    if data.get("description"):
        _description_block(slide, data["description"], Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.6))
    _textbox(slide, "KEY POINTS", Inches(0.5), Inches(2.3), Inches(5.8), Inches(0.4), 10, THEME["col_label"], bold=True)
    _textbox(slide, "DETAILS",    Inches(6.8), Inches(2.3), Inches(5.8), Inches(0.4), 10, THEME["col_label"], bold=True)
    _bullets_block(slide, left_b,  Inches(0.5), Inches(2.8), Inches(5.8), Inches(4.0))
    _bullets_block(slide, right_b, Inches(6.8), Inches(2.8), Inches(5.8), Inches(4.0))

def _render_quote(slide, data):
    _textbox(slide, f'"{data.get("quote", "")}"',
             Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.0),
             22, THEME["title_color"], align=PP_ALIGN.CENTER)
    _textbox(slide, f"— {data.get('author', '')}",
             Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6),
             15, THEME["accent"], bold=True, align=PP_ALIGN.CENTER)
    if data.get("description"):
        _description_block(slide, data["description"], Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.9))

def _add_stat_card(slide, value, label, x, y, w):
    # card background box matching .stat-card (#f7f5ff with border)
    card = slide.shapes.add_shape(1, x, y, w - Inches(0.2), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = THEME["stat_bg"]
    card.line.color.rgb = RGBColor(0xE0, 0xD9, 0xFF)
    _textbox(slide, value, x + Inches(0.1), y + Inches(0.1), w - Inches(0.4), Inches(0.7),
             32, THEME["stat_color"], bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, label.upper(), x + Inches(0.1), y + Inches(0.85), w - Inches(0.4), Inches(0.4),
             10, THEME["muted"], align=PP_ALIGN.CENTER)

def _render_stats(slide, data):
    stats = data.get("stats", [])
    if data.get("description"):
        _description_block(slide, data["description"], Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.6))
    if stats:
        col_w = Inches(12.3) / len(stats)
        for i, stat in enumerate(stats):
            _add_stat_card(slide, stat.get("value", ""), stat.get("label", ""),
                           Inches(0.5) + col_w * i, Inches(2.3), col_w)
    if data.get("bullets"):
        _bullets_block(slide, data["bullets"], Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8), font_size=14)

# ── MCP Server ────────────────────────────────────────────────────────────────
server = Server("napkin-ppt-mcp-server")

@server.list_tools()
async def list_tools() -> list[types.Tool]:
    return [
        types.Tool(
            name="create_presentation",
            description="Initialize a new blank presentation. Must be called first.",
            inputSchema={
                "type": "object",
                "properties": {
                    "output_path": {"type": "string", "description": "File path for the .pptx output"},
                },
                "required": ["output_path"],
            },
        ),
        types.Tool(
            name="add_title_slide",
            description="Add the opening title slide with title, subtitle, and optional cover image.",
            inputSchema={
                "type": "object",
                "properties": {
                    "title":        {"type": "string"},
                    "subtitle":     {"type": "string"},
                    "image_url":    {"type": "string", "description": "Direct URL to an image to embed as cover"},
                },
                "required": ["title"],
            },
        ),
        types.Tool(
            name="add_slide",
            description=(
                "Add a content slide. layout must be one of: bullets, two_column, quote, stats. "
                "bullets: needs 'bullets' list. "
                "two_column: needs 'left_bullets' and 'right_bullets'. "
                "quote: needs 'quote' and 'author'. "
                "stats: needs 'stats' list of {value, label} and optional 'bullets'."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "title":         {"type": "string"},
                    "layout":        {"type": "string", "enum": ["bullets", "two_column", "quote", "stats"]},
                    "bullets":       {"type": "array",  "items": {"type": "string"}},
                    "left_bullets":  {"type": "array",  "items": {"type": "string"}},
                    "right_bullets": {"type": "array",  "items": {"type": "string"}},
                    "quote":         {"type": "string"},
                    "author":        {"type": "string"},
                    "stats":         {"type": "array",  "items": {"type": "object"}},
                    "description":   {"type": "string"},
                    "image_url":     {"type": "string"},
                    "notes":         {"type": "string"},
                },
                "required": ["title", "layout"],
            },
        ),
        types.Tool(
            name="save_presentation",
            description="Write the presentation to disk. Call this last.",
            inputSchema={"type": "object", "properties": {}},
        ),
    ]

@server.call_tool()
async def call_tool(name: str, arguments: dict) -> list[types.TextContent]:
    global _prs, _output_path, _used_image_urls

    if name == "create_presentation":
        _prs = PRS()
        _prs.slide_width  = Inches(13.33)
        _prs.slide_height = Inches(7.5)
        _output_path = arguments["output_path"]
        _used_image_urls = set()
        return [types.TextContent(type="text", text=f"Presentation created → '{_output_path}'.")]

    if name == "add_title_slide":
        if _prs is None:
            return [types.TextContent(type="text", text="ERROR: call create_presentation first.")]
        ts = _prs.slides.add_slide(_prs.slide_layouts[6])
        if arguments.get("image_url"):
            # image on right half, white bg on left — matches TitleSlide CSS layout
            _set_bg(ts, THEME["bg"])
            _add_image(ts, arguments["image_url"], Inches(6.5), Inches(0), Inches(6.83), Inches(7.5))
        else:
            _set_bg(ts, THEME["bg"])
        # left content panel
        _textbox(ts, arguments["title"], Inches(0.7), Inches(2.0), Inches(5.5), Inches(2.0), 38, THEME["title_color"], bold=True)
        if arguments.get("subtitle"):
            _textbox(ts, arguments["subtitle"], Inches(0.7), Inches(4.2), Inches(5.5), Inches(0.8), 17, THEME["muted"])
        _textbox(ts, "Generated by Napkin AI", Inches(0.7), Inches(5.2), Inches(5.5), Inches(0.4), 12, THEME["accent"])
        return [types.TextContent(type="text", text=f"Title slide '{arguments['title']}' added.")]

    if name == "add_slide":
        if _prs is None:
            return [types.TextContent(type="text", text="ERROR: call create_presentation first.")]
        slide = _prs.slides.add_slide(_prs.slide_layouts[6])
        _set_bg(slide)
        layout = arguments.get("layout", "bullets")

        _textbox(slide, arguments["title"], Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9), 26, THEME["title_color"], bold=True)
        _divider(slide, Inches(1.2))

        if layout == "two_column":
            _render_two_column(slide, arguments)
        elif layout == "quote":
            _render_quote(slide, arguments)
        elif layout == "stats":
            _render_stats(slide, arguments)
        else:
            _render_bullets(slide, arguments)

        if arguments.get("notes"):
            slide.notes_slide.notes_text_frame.text = arguments["notes"]

        return [types.TextContent(type="text", text=f"Slide '{arguments['title']}' ({layout}) added.")]

    if name == "save_presentation":
        if _prs is None:
            return [types.TextContent(type="text", text="ERROR: no presentation to save.")]
        buf = io.BytesIO()
        _prs.save(buf)
        with open(_output_path, "wb") as f:
            f.write(buf.getvalue())
        return [types.TextContent(type="text", text=f"Saved '{_output_path}' ({len(_prs.slides)} slides).")]

    return [types.TextContent(type="text", text=f"Unknown tool: {name}")]


async def main():
    async with stdio_server() as (r, w):
        await server.run(r, w, server.create_initialization_options())

if __name__ == "__main__":
    asyncio.run(main())
